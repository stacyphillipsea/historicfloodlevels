import requests
import json
import pandas as pd
import dash
from dash import dash_table
from dash.dependencies import Input, Output
import dash_core_components as dcc
import dash_html_components as html
import dash_bootstrap_components as dbc
from datetime import datetime
import plotly.graph_objects as go
import numpy as np
import pandas as pd
import os
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Pt, Inches
import folium


### GET YOUR DATA BITS
# Define key constants
BASE_URL = "http://environment.data.gov.uk/hydrology/id"
BASE_STATIONS_URL = "http://environment.data.gov.uk/hydrology/id/stations"
WISKI_IDS = ['2175', '2077', '2134', '2180', '2001', '2642', '2085', '2616', '2032', '2087', '2606', '2057', '2071',
             '2618', '2165', '2102', '2153', '2132', '2008', '2086', '2002', '2128', '055829', '055811', '055807',
             '055817', '055843', '4143', '4078', '4018', '4703', '4052', '4040', '4083', '4006', '4012', '4019',
             '4069', '4039', '4066', '4081', '4878', '4003', '2090', '2019', '2091', '2093', '2050', '2049', '2048',
             '452039', '2092', '2621', '2104', '2531', '055041', '055003', '055816']
# ## WYE WISKI IDS
# WISKI_IDS =['055002', '055028', '055811', '055040', '055021','055014', '055829', '055041', '055003', '055843',
#             '055843', '055807', '055039', '055817', '055031', '055013', '055018']
#WISKI_IDS = ['2175', '2077', '2134']
MIN_DATE_STR = "2023-10-01"
MAX_DATE_STR = "2024-02-29"
MIN_DATE = datetime.strptime(MIN_DATE_STR, '%Y-%m-%d')
MAX_DATE = datetime.strptime(MAX_DATE_STR, '%Y-%m-%d')
DATE_FILTERS = {
    'Babet': ('2023-10-18', '2023-10-31', 'red'),
    'Ciaran': ('2023-11-01', '2023-11-08', 'blue'),
    'Elin & Fergus': ('2023-11-09', '2023-12-17', 'black'),
    'Gerrit': ('2023-12-26', '2024-01-02', 'pink'),
    'Henk': ('2024-01-02', '2024-01-11', 'green'),
    'Isha & Jocelyn': ('2024-01-21', '2024-01-27', 'orange'),
    'Early February': ('2024-02-09', '2024-02-12', 'cornflowerblue'),
    'Late February': ('2024-02-22', '2024-02-25', 'gray')
}

# Load data
sites_of_interest_merge = pd.read_csv('sites_of_interest_merge.csv')
gaugeboard_data = pd.read_csv('gaugeboard_data.csv')
threshold_values = sites_of_interest_merge[sites_of_interest_merge['Threshold'].notnull()]
threshold_values['Threshold'] = threshold_values['Threshold'].astype(float)
threshold_dict = threshold_values.set_index('Gauge')['Threshold'].to_dict()
ea_logo = "EA_logo.jpg"
ea_logo_clip = "https://png2.cleanpng.com/sh/21bbe3c8fc2fbad88acec34e033dd3de/L0KzQYm3VMIxN5N0fZH0aYP2gLBuTfVvfpp3h9D2ZX73PbLuhf5kgV5teexqcnTyhcS0lBF0fJYyhtN9dYLkfH7sjwZqepDzRdd3dnn1f7B0hf51aZ0yhtN9dYLoPYbohMllP5MASNRrNEW1Poa5V8k4OmM6Sac7NEK1RYqAV8A1QF91htk=/kisspng-environment-agency-hazardous-waste-natural-environ-environmental-nature-5ad9d7b90bb452.527972251524225977048.png"
dash_logo = "https://raw.githubusercontent.com/tomkdefra/defra-dash/main/assets/images/DASH_logo_colour.png"

## Creat navigation bar/header
NAVBAR = dbc.Navbar(
    dbc.Container(
        [
            html.A(
                dbc.Row(
                    [
                        dbc.Col(html.Img(src="https://visualpharm.com/assets/818/Cheese-595b40b65ba036ed117d2a6c.svg", 
                              height="50px")),
                        dbc.Col(dbc.NavbarBrand("FETA Homepage", className="ms-2", 
                                                style={"color":"#008531", "font-weight": "bold"})),
                    ],
                    align="center",
                    className="g-0",
                ),
                href="https://dap-prd2-connect.azure.defra.cloud/level_data/",
                style={"textDecoration": "none"},
            ),
            dbc.Row(
                [
                dbc.Col(dbc.NavbarBrand("App powered by", className="ms-2", 
                                        style={"color":"#008531"}), width="auto"),
                dbc.Col(html.A(html.Img(src=ea_logo_clip, height="60px")), width="auto"),
                dbc.Col(html.A(html.Img(src=dash_logo, height="50px")), width="auto"),
                ],
                className="g-0 ms-auto flex-nowrap mt-3 mt-md-0", # no gutter, right hand side, no wrap, margins
                align="center",
                justify="end"  # Align the row content to the end (right-hand side)
            ),
        ]
    ),
    color="#d9f5ce",
    dark=True,
)

### MAKE YOUR FUNCTIONS
# Fetch data for a single station
def fetch_station_data(wiski_id):
    try:
        url_endpoint = f"{BASE_STATIONS_URL}?wiskiID={wiski_id}"
        response = requests.get(url_endpoint)
        response.raise_for_status()
        data = json.loads(response.content)
        if 'items' in data and data['items']:
            label_field = data['items'][0].get('label')
            name = str(label_field[1] if isinstance(label_field, list) else label_field)
            river_name = data['items'][0].get('riverName')
            latitude = data['items'][0].get('lat')
            longitude = data['items'][0].get('long')
            measure_url = f"{BASE_URL}/measures?station.wiskiID={wiski_id}&observedProperty=waterLevel&periodName=15min"
            response = requests.get(measure_url)
            response.raise_for_status()
            measure = json.loads(response.content)
            if 'items' in measure and measure['items']:
                measure_id = measure['items'][0]['@id']
                readings_url = f"{measure_id}/readings?mineq-date={MIN_DATE_STR}&maxeq-date={MAX_DATE_STR}"
                response = requests.get(readings_url)
                response.raise_for_status()
                readings = json.loads(response.content)
                readings_items = readings.get('items', [])
                if readings_items:
                    df = pd.DataFrame.from_dict(readings_items)
                    df['dateTime'] = pd.to_datetime(df['dateTime'])
                    return {
                        'name': name,
                        'date_values': df[['dateTime', 'value']],
                        'river_name': river_name,
                        'lat': latitude,
                        'long': longitude
                    }
                else:
                    print(f"No readings found for {name} ({wiski_id})")
            else:
                print(f"No measure items found for WISKI ID {wiski_id}")
        else:
            print(f"No station items found for WISKI ID {wiski_id}")
    except requests.exceptions.RequestException as e:
        print(f"Error fetching data for WISKI ID {wiski_id}: {e}")
    return None

# Fetch data for all stations
def fetch_all_station_data():
    data_dict = {}
    for wiski_id in WISKI_IDS:
        station_data = fetch_station_data(wiski_id)
        if station_data:
            data_dict[station_data['name']] = station_data
    return data_dict

# Find maximum values for each filter
def find_max_values(df, filters):
    max_values = {}
    for filter_name, date_range in filters.items():
        min_date, max_date, color = date_range
        condition = (df['dateTime'] >= min_date) & (df['dateTime'] <= max_date)
        filtered_df = df[condition].dropna()  # Drop rows with NaN values
        if not filtered_df.empty:
            max_value_row = filtered_df.loc[filtered_df['value'].idxmax(), ['dateTime', 'value']]
            max_value_row['value'] = round(max_value_row['value'], 2)  # Round the maximum value to 2 decimal places
            max_values[filter_name] = max_value_row
    return max_values

# Find and store maximum values for all stations
def find_and_store_max_values(data_dict):
    max_values = {}
    for station_name, station_data in data_dict.items():
        df = station_data.get('date_values')
        if df is not None:
            max_values[station_name] = find_max_values(df, DATE_FILTERS)
    return max_values

#Generate storm info for the peak table
def generate_storm_info():
    storm_info = []
    for storm, (start_date, end_date, _) in DATE_FILTERS.items():
        formatted_start_date = datetime.strptime(start_date, '%Y-%m-%d').strftime('%d %b %Y')
        formatted_end_date = datetime.strptime(end_date, '%Y-%m-%d').strftime('%d %b %Y')
        storm_info.append(html.Div([
            html.P(f"{storm}: {formatted_start_date} to {formatted_end_date}", style={'font-size': '14px'}),
        ]))
    return storm_info

# Create a table for all the sites
def process_peak_table_all(max_values, sites_of_interest_merge):
    # Create an empty DataFrame
    df_list = []

    # Iterate through the outer dictionary
    for station, inner_dict in max_values.items():
        # Iterate through the inner dictionary
        for storm, values in inner_dict.items():
            # Extract dateTime and value from the Series object
            if values is not None:
                date_time = values.get('dateTime')
                value = values.get('value')
            else:
                date_time, value = None, None

            # Create a new row as a dictionary
            row_dict = {'Station': station, 'Storm': storm, 'DateTime': date_time, 'Value': value}
            # Append the row dictionary to the list
            df_list.append(row_dict)

    # Create the DataFrame
    df = pd.DataFrame(df_list)

    # Pivot the DataFrame
    pivot_df = df.pivot_table(index='Station', columns='Storm', values=['DateTime', 'Value'], aggfunc='first')

    # Flatten the multi-level columns
    pivot_df.columns = [f'{col[0]}_{col[1]}' for col in pivot_df.columns]

    # Create a new DataFrame with the flat structure
    flat_df = pd.DataFrame()

    # Iterate through the original DataFrame columns and create flat structure
    for storm in df['Storm'].unique():
        flat_df[f'{storm}_DateTime'] = pivot_df[f'DateTime_{storm}'].dt.strftime('%d-%b-%Y %H:%M')
        flat_df[f'{storm}_Value'] = pivot_df[f'Value_{storm}']

    # Reset index if necessary
    flat_df.reset_index(inplace=True)

# Merge with 'sites_of_interest_merge' DataFrame
    peak_table_all = pd.merge(flat_df, sites_of_interest_merge[['Region', 'River','Gauge','Order']], left_on='Station', right_on='Gauge')
    columns_to_move = ['Order','Region', 'River']
    new_order = columns_to_move + [col for col in peak_table_all.columns if col not in columns_to_move]
    peak_table_all = peak_table_all[new_order]
    peak_table_all.drop(columns=['Gauge'], inplace=True)

    peak_table_all = peak_table_all.sort_values(by=['Order'], ascending=[True])
    peak_table_all.set_index('Order', inplace=True)

    peak_table_all = peak_table_all.drop_duplicates(subset=['Station'])
    # peak_table_all = flat_df
    
    return df, peak_table_all

# Do the comparison to the gaugeboard data
def gaugeboard_comparison(gaugeboard_data, df):
    # Format gaugeboard data datetimes
    gaugeboard_data['Date'] = pd.to_datetime(gaugeboard_data['Date'], format='%d/%m/%Y').dt.date

    # Rename 'Gauge' column to 'Station' in gaugeboard
    gaugeboard_data.rename(columns={'Gauge': 'Station'}, inplace=True)

    # Add a new column 'Storm' filled with nulls in gaugeboard
    gaugeboard_data['Storm'] = None

    # Change 'datetime' column to 'date' in df
    df['DateTime'] = df['DateTime'].dt.date
    df.rename(columns={'Value': 'Level'}, inplace=True)
    df.rename(columns={'DateTime': 'Date'}, inplace=True)

    comparison_concat = pd.concat([gaugeboard_data, df], ignore_index=True)

    comparison_concat['Station'] = comparison_concat['Station'].replace({
        'Bewdley us': 'Bewdley',
        'SAXONS LODE US': 'Saxons Lode',
        'Buildwas Us': 'Buildwas'
    })

    # Group the DataFrame by 'station' and sort each group by 'level' from highest to lowest
    sorted_df = comparison_concat.groupby('Station', as_index=False).apply(lambda x: x.sort_values(by='Level', ascending=False))

    # Add a new column 'Ranking' for each group indicating the rank of each level
    sorted_df['Ranking'] = sorted_df.groupby('Station').cumcount() + 1

    # Add a new column for difference from peak
    sorted_df['Difference_from_peak'] = sorted_df.groupby('Station')['Level'].transform(lambda x: round(x.max() - x, 2))

    # Convert 'Date' column to datetime format after sorting
    sorted_df['Date'] = pd.to_datetime(sorted_df['Date'])

    # Add a new column for difference from peak
    sorted_df['Years_since_peak'] = sorted_df.groupby('Station').apply(
        lambda group: round(abs((group['Date'] - group.loc[group['Level'].idxmax(), 'Date']).dt.days / 365.25))
        ).reset_index(level=0, drop=True)
    
    # Reset the index to flatten the DataFrame
    sorted_df.reset_index(drop=True, inplace=True)

    # Display the sorted DataFrame
    sorted_df.head(20)

    # Count the number of rows for each unique value in 'Station' column
    station_counts = sorted_df['Station'].value_counts()

    # Filter out stations with 5 or fewer records
    stations_to_keep = station_counts[station_counts > 8].index

    # Filter the DataFrame to keep only the stations with more than 5 records
    filtered_df = sorted_df[sorted_df['Station'].isin(stations_to_keep)]

    # Display the filtered DataFrame
    filtered_df['Station'].unique()

    ranked_df = filtered_df[filtered_df['Storm'].notna()]
    ranked_df.head(20)

    top_ten = ranked_df[ranked_df['Ranking'] <= 10]
    top_ten = top_ten.sort_values(by='Ranking')
    top_ten['Date'] = pd.to_datetime(top_ten['Date'], format='%d/%m/%Y')  
    top_ten['Date'] = top_ten['Date'].dt.strftime('%d-%b-%Y')
    
    return top_ten, filtered_df

# Make a top 10 list for the station selected
def station_top_ten(selected_station):
    if selected_station:
        # Filter for the selected station
        station_df = filtered_df[filtered_df["Station"] == selected_station]
        # Get the top 10 values for that station
        top_ten_df = station_df.nlargest(n=10, columns= 'Level')
        top_ten_df['Date'] = pd.to_datetime(top_ten_df['Date'], format='%d/%m/%Y')  
        top_ten_df['Date'] = top_ten_df['Date'].dt.strftime('%d-%b-%Y')
        return top_ten_df
    else:
        return pd.DataFrame()

# Make the historic level plots
def plot_historic_levels(filtered_df, selected_station, threshold_dict):
    # Check if 'Station' column is present in filtered_df
    if 'Station' not in filtered_df.columns:
        # Return an empty figure if 'Station' column is not present
        return go.Figure()
    
    # Check if the selected station is valid
    if selected_station and selected_station in filtered_df['Station'].unique():
        # Filter the DataFrame for the specified station name
        station_df = filtered_df[filtered_df["Station"] == selected_station]
        
        # Filter the DataFrame for top 3 values
        top_3_values = station_df.nlargest(n=3, columns='Level')
        
        # Filter the DataFrame for rows where 'Storm' is not None
        storm_not_none = station_df[station_df["Storm"].notnull()]
        
        # Concatenate the two DataFrames
        result_df = pd.concat([top_3_values, storm_not_none])

        # Determine the maximum level in the data
        max_level = result_df['Level'].max()
        
        # Round up to the nearest whole number
        max_level = np.ceil(max_level)
        
        # Create a new figure
        fig = go.Figure()
        
        # Iterate over each row in the filtered DataFrame
        for index, row in result_df.iterrows():
            # Assign colors based on Ranking
            if row['Ranking'] == 1:
                color = 'red'
            elif row['Ranking'] == 2:
                color = 'orange'
            elif row['Ranking'] == 3:
                color = 'yellow'
            else:
                color = 'blue'  # Default color
            
            # Format the date string
            formatted_date = row['Date'].strftime('%d-%b-%Y')
            
            # Add a trace for each row
            fig.add_trace(go.Scatter(x=[1, 2], y=[row['Level'], row['Level']], mode='lines+markers', line=dict(color=color), 
                                     name=f"{formatted_date} {row['Storm'] if pd.notnull(row['Storm']) else ''}",
                                     text=[f"Date: {formatted_date}<br>Level: {row['Level']}<br>Storm: {row['Storm'] if pd.notnull(row['Storm']) else 'Historic'}"] * 2,
                                     hoverinfo='text'))
        
        # Add manual trace for the line at typical max range
        if selected_station in threshold_dict:
            typical_max_range = threshold_dict[selected_station]
            fig.add_trace(go.Scatter(x=[1, 2], y=[typical_max_range, typical_max_range], mode='lines+markers', line=dict(color='lightskyblue'), 
                                     name='Typical max range', text=f"Typical max range: {typical_max_range}", hoverinfo='text'))
        
        # Update layout...
        fig.update_layout(
            xaxis=dict(
                showticklabels=False,
                range=[0.5, 2.5]
            ),
            yaxis=dict(
                title='Level (above gaugeboard datum)',
                rangemode='tozero',  # Ensure y-axis starts at zero
                tickmode='linear',   # Use linear ticks
                tick0=0,             # Start tick at zero
                dtick=1,             # Set tick interval to 1
                tickvals=np.arange(0, max_level + 1),  # Set tick values to whole numbers up to max_level
                range=[0, max_level]  # Set y-axis range to start from 0 and end at max_level
            ),
            title={
                'text': f"Historic levels for {selected_station}",
                'x': 0.5,  # Set x position to center
                'y': 0.95   # Adjust y position as needed
            },
            legend=dict(
                title='Recorded levels and dates',
                x=1.2,  # Increase distance from right edge
            ),
            width=400,  # Adjust overall width
            plot_bgcolor='white'
        )
        
        fig.update_yaxes(
            mirror=True,
            ticks='outside',
            showline=True,
            linecolor='black',
            gridcolor='lightgrey'
        )
        
        return fig
    else:
        # If the selected station is not valid or no station is selected, return an empty figure
        return go.Figure()

# Function to load station data from JSON file
def load_station_data_from_json(file_path):
    try:
        # Load data from JSON file
        with open(file_path, "r") as json_file:
            data_dict = json.load(json_file)
        
        # Convert date_values from JSON strings to DataFrames
        for station_data in data_dict.values():
            station_data['date_values'] = pd.read_json(station_data['date_values'])

        return data_dict
    except FileNotFoundError:
        print(f"File {file_path} not found.")
        return None

# Define a list of colors to use as a palette
color_palette = [
    'blue', 'green', 'cadetblue', 'orange', 'purple', 'pink', 'gray', 'beige', 
    'lightblue', 'red', 'lightred', 'lightgray', 'darkblue', 'darkgreen', 'darkred', 'darkpurple'
]

# Function to create Folium map with markers for stations
def create_map(data_dict, selected_station=None):
    # Create Folium map centred on Hagley (roughly in centre of WMD)
    m = folium.Map(location=[52.4083, -2.2272], zoom_start=10)
    
    # Extract unique river names from the data_dict
    unique_rivers = sorted(set(station_data.get('river_name', None) for station_data in data_dict.values() if 'river_name' in station_data))

    # Create river-color mapping by assigning colors from the palette
    river_color_mapping = {river: color_palette[i % len(color_palette)] for i, river in enumerate(unique_rivers)}

    # Add markers for all stations
    for station_name, station_data in data_dict.items():
        lat = station_data.get('lat', None)
        long = station_data.get('long', None)
        river_name = station_data.get('river_name', None)
        popup_content = f"<b>{station_name}</b><br>{river_name}"
        
        if lat is not None and long is not None:
            # Select marker color based on river name
            marker_color = river_color_mapping.get(river_name, 'gray')  # Default to gray if river name not found
            # Add marker for station with selected color
            folium.Marker(location=[lat, long], popup=popup_content, 
                          icon=folium.Icon(icon = "info-sign", color=marker_color, icon_color="white")).add_to(m)
            m.get_root().header.add_child(folium.Element("<style>.leaflet-popup-content { width: 100px; text-align: center; }</style>"))

    # Adjust zoom level and center map if a station is selected
    if selected_station:
        station_data = data_dict[selected_station]
        lat = station_data.get('lat', None)
        long = station_data.get('long', None)

        if lat is not None and long is not None:
            # If latitude and longitude are available for the selected station, center the map on that station
            m.location = [lat, long]
            m.zoom_start = 10  # Adjust the zoom level as needed

    # Convert Folium map to HTML
    map_html = m.get_root().render()

    return map_html

#### FUNCTION TO MAKE DICTIONARY OFFLINE AND THEN LOAD

# # Fetch data for all stations
# data_dict = fetch_all_station_data()

# def fetch_and_save_all_station_data():
#     data_dict = {}
#     for wiski_id in WISKI_IDS:
#         station_data = fetch_station_data(wiski_id)
#         if station_data:
#             # Convert DataFrame to JSON-serializable format
#             date_values_json = station_data['date_values'].to_json(orient='records')
#             # Replace DataFrame with JSON string in station_data dictionary
#             station_data['date_values'] = date_values_json
#             data_dict[station_data['name']] = station_data

#     # Specify the file path where you want to save the JSON file
#     file_path = "C:\\Users\\SPHILLIPS03\\Documents\\repos\\levels_app_folder\\nested_dict.json"

#     # Save the dictionary containing station data to a JSON file
#     with open(file_path, "w") as json_file:
#         json.dump(data_dict, json_file)

#     print("JSON file saved successfully.")

# fetch_and_save_all_station_data()

### CALL YOUR FUNCTIONS 
# Load station data from JSON file
file_path = "nested_dict.json"
data_dict = load_station_data_from_json(file_path)

if data_dict:
    print("Data loaded successfully.")
else:
    print("Error loading data.")

# Find and store maximum values for all stations
max_values = find_and_store_max_values(data_dict)

# Create peak table DataFrame
df, peak_table_all = process_peak_table_all(max_values, sites_of_interest_merge)

# Call gaugeboard_comparison function
top_ten_records, filtered_df = gaugeboard_comparison(gaugeboard_data, df)

# Identify the common station that has historic values, threholds and peak data
complete_stations = sorted(set(filtered_df['Station'].unique()) & set(threshold_dict.keys()) & set(data_dict.keys()))
percent_complete = len(complete_stations) / len(data_dict) * 100 if len(data_dict) > 0 else 0

initial_map_html = create_map(data_dict)

#### SAVING STUFF FOR POWERPOINT PRESENTATION
# Save all the charts for use in the Powerpoint
import matplotlib.pyplot as plt
def save_all_station_charts(data_dict, output_directory):
    # Create the output directory if it doesn't exist
    if not os.path.exists(output_directory):
        os.makedirs(output_directory)

    image_paths = []  # Initialize list to store image paths
    
    for station_name, station_data in data_dict.items():
        df = station_data.get('date_values')
        if df is not None:
            # Create the figure
            plt.figure(figsize=(10, 3))
            plt.plot(df['dateTime'], df['value'])
            # plt.title(f'River Levels for {station_name} ({station_data.get("river_name", "Unknown River")})')
            plt.xlabel('Date Time')
            plt.ylabel('Value')
            
            # Plot max values if available
            if station_name in max_values:
                for filter_name, max_value_info in max_values[station_name].items():
                    max_datetime = max_value_info['dateTime']
                    max_value = max_value_info['value']
                    color = DATE_FILTERS[filter_name][2]
                    plt.plot(max_datetime, max_value, marker='o', markersize=10, color=color, label=f'Storm {filter_name} peak')
            
            # Plot legend outside plot
            plt.legend(bbox_to_anchor=(1.05, 1), loc='upper left')  # Add legend

            # Get river name from data_dict
            river_name = data_dict.get(station_name, {}).get('river_name', 'Unknown River')

            # Specify the path to save the image
            image_path = os.path.join(output_directory, f"peakchart_{river_name}_{station_name}.png")
            
            # Save the figure as an image
            plt.tight_layout()  # Adjust layout to prevent clipping
            plt.savefig(image_path, bbox_inches='tight')  # Use bbox_inches='tight' to include legend
            plt.close()  # Close the figure to free up memory

            # Append the image path to the list
            image_paths.append(image_path)

    return image_paths  # Return list of image paths

output_directory = "C:\\Users\\SPHILLIPS03\\Documents\\repos\\levels_app_folder_charts"
image_paths = save_all_station_charts(data_dict, output_directory)

def insert_plots_into_ppt(image_paths, prs):
    for image_path in image_paths:
        # Extract river name and location from the filename
        file_name = os.path.splitext(os.path.basename(image_path))[0]
        parts = file_name.split('_')
        river_name = ' '.join(parts[1:-1])  #River name is everything between first and last underscore  
        station = parts[-1]  

        # Insert each image onto a new slide in the presentation
        slide = prs.slides.add_slide(prs.slide_layouts[5])  # Slide layout 5
    
        # # Modify title text
        title_text = f"River Levels for {station}\n({river_name})"
        title_shape = slide.shapes.title
        title_shape.text = title_text
        
        # Set font size for each line
        title_shape.text_frame.paragraphs[0].font.size = Pt(36)  # Font size for station
        title_shape.text_frame.paragraphs[1].font.size = Pt(24)  # Font size for river name
        
        # Set font bold for river name
        title_shape.text_frame.paragraphs[1].font.bold = True
        
        # Set vertical alignment for the title text
        # title_shape.text_frame.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        
        # Set image dimensions and position
        image_width = Inches(10)  # Adjust as needed
        image_height = Inches(3)  # Adjust as needed
        image_left = (prs.slide_width - image_width) / 2
        image_top = Inches(1.5)
        
        # Add the plot image to the slide
        slide.shapes.add_picture(image_path, image_left, image_top, image_width, image_height)

# Create a PowerPoint presentation object
prs = Presentation()

# Call the function to insert plots into the presentation
insert_plots_into_ppt(image_paths, prs)

# Save the PowerPoint presentation
output_presentation_path = "C:\\Users\\SPHILLIPS03\\Documents\\repos\\levels_app_folder_exports\Winter2324_PeakPlots.pptx"
prs.save(output_presentation_path)


### MAKE YOUR APP
# Initialize Dash app
external_stylesheets = [dbc.themes.MINTY]
app = dash.Dash('__main__', external_stylesheets=external_stylesheets)

# Define the content of the modal
modal_content = dbc.Modal([
    dbc.ModalHeader("A note on data completeness"),
    dbc.ModalBody([
        html.H4("Not all stations have historic values, typical ranges, and peak values", style={"textAlign": "left", "color": "green", "fontWeight": "bold"}),
        html.P("Stations with complete datasets:", style={"textAlign": "left", "font-size": "16px", "color": "green"}),
        html.P(', '.join(complete_stations), style={"textAlign": "left", "font-size": "14px", "color": "green"}),
        html.P(f"That is {len(complete_stations)} stations out of {len(data_dict)} in the whole dataset ({percent_complete:.0f}%)", style={"textAlign": "left", "font-size": "12px", "color": "green", "fontStyle": "italic"}),
        dbc.CardImg(src="https://media2.giphy.com/media/1Zp8tUAMkOZDMkqcHb/giphy.gif?cid=6c09b952rjrtfs3brpsa0z89g2oeqrzgg7d3sdoj8fon3aqd&ep=v1_internal_gif_by_id&rid=giphy.gif&ct=g", bottom=True, style={"width": "250px"}), 
    ]),
    dbc.ModalFooter(
        dbc.Button("Close", id="close_modal", className="ml-auto")
    ),
], id="modal")


### DEFINE APP LAYOUT
app.layout = dbc.Container([
    NAVBAR,
    # Titles
    html.H1("Welcome to the Flood Event Telemetry Analyser (FETA)!", style={"textAlign":"center", }),  # title
    html.H5(["This app allows you to explore river level data for the Winter flood events 2023-2024",
            html.Br(),
            "for sites across the West Midlands"],
            style={"textAlign":"center"}), 
    html.H4(["App created by ", html.A("Stacy Phillips", href="mailto:stacy.phillips1@environment-agency.gov.uk?subject=I%20love%20your%20FETA%20app!")], 
            style={"textAlign": "center", "color":"#034B89"}),         
    html.Hr(),  # line break
    
    # How to use information & photo
    html.Div([
        dbc.Row([
            dbc.Col([
                html.H5("How to use the FETA", style={"text-decoration": "underline"}),
                html.P("Data is automatically downloaded using the Environment Agency API for a number of different specified gauging stations."),
                html.P(["Choose your station of interest ",
                        html.A("(here)", href="#station-choice"),
                        " and then look at the peak levels for that station across different storm events ",
                        html.A("(here)", href="#peak-info"),
                        "."]),
                html.P(["You can then compare these peak levels to historic records ",
                        html.A("(here)", href="#historic-info"),
                        " (where available: ",
                        html.A("see note here", href="#data-modal"), 
                        ") to give context for each storm event."]),     
                html.P("All of the charts, maps and tables are interactive in some way, allowing you to filter, sort and investigate the data as you please."),
                html.P(["Further information about data sources can be found at the ",
                        html.A("bottom of the page", href="#data-info"),
                        "."]),
                html.Div(style={"height": "14px"}),
            ], width=9),
            dbc.Col([
                html.Img(src="https://www.shropshirestar.com/resizer/R-JBYJySB7d1sV88kojEsofoO0w=/1200x675/cloudfront-us-east-1.images.arcpublishing.com/mna/6RGV7KRRMRDC7JMDGNQFHD7X7U.jpg", style={"width": "100%"}),
                html.P("Flooding in Shrewsbury following Storm Gerrit", 
                       style={"textAlign": "center", "fontStyle": "italic", "color": "gray", "font-size":"12px", "margin-bottom": "5px"}),
                html.P("Photo from Shropshire Star", 
                       style={"textAlign": "center", "fontStyle": "italic", "color": "gray", "font-size":"10px"})  # Adjust margin-top
            ]),
        ]),
        dbc.Row([
            dbc.Col(
                html.H6("!!!This site is a work in progress: Stiltonnes of work to do!!!", 
                        style={"textAlign":"center", "color": "red", "fontStyle": "italic", "fontWeight": "bold"}),
            )
        ]),
    ]),       
    html.Hr(),  # line break
    
    # Data info modal popup
    html.Div([
        html.Div(id="data-modal"),
        html.Button("A note on data completeness", id="open_modal"),
        # Define the modal
        modal_content,
    ]),
    html.Div(style={'height': '10px'}),  # Adding vertical space with a div
    # Dropdowns and map
    dbc.Row([
        dbc.Col([
            html.P("Choose a site for peak analysis:", id="station-choice", 
                   style={'font-size': '16px', "fontStyle": 'bold'}), 
            html.P("First, pick the river name, and then the stations available for that river will be shown.", 
                   style={'font-size': '14px'}),
            html.P("You can start typing into the bar to search, or pick from the dropdown.", 
                   style={'font-size': '14px'}),
            html.P("The map to the right can help you identify what river a site is on; each river is coloured differently.", 
                   style={'font-size': '14px'}),
            dcc.Dropdown(
                id="river-dropdown",
                clearable=False,
                value="River Severn",  # Default value for the river dropdown
                options=[
                    {'label': river_name, 'value': river_name} for river_name in sorted(set([v['river_name'] for v in data_dict.values()]))
                ],
                style={'font-size': '16px'}
            ),
            dcc.Dropdown(
                id="station-dropdown",
                clearable=False,
                value="Welsh Bridge",  # Default value for the station dropdown
                options=[],  # Will be dynamically populated based on the selected river name
                style={'font-size': '16px'}
            ),
        ], width=8),
        dbc.Col(
            html.Iframe(id='map-container',srcDoc=initial_map_html, width='100%', height='200'),
            width=4
        ),
    ]),
    html.Hr(),  # line break
    # Peak chart and peak table
    dbc.Row([
        dbc.Col(
            html.Div([
                html.H2("River levels for this station over Winter 23-24",id="peak-info", 
                        style={"textAlign": "center", 'font-size': '14px'}),
                html.Div(id="output-graph", className="card"),
            ]),
            width=8
        ),  
        dbc.Col(
            html.Div([
                html.H2("Peaks identified for this station over Winter 23-24", style={"textAlign": "center", 'font-size': '14px'}),
                html.Div(id="peak-table", className="card"),
            ]),
            width=4
        ),  
    ]),
    html.Hr(),  # line break
    # Peaks versus historic table and chart
    dbc.Row([
        dbc.Col(
            html.Div([
                html.H2("Top 10 historic levels for this station", id="historic-info", style={"textAlign": "center", 'font-size': '14px'}),
                dash_table.DataTable(
                    id='station-top-ten-table',
                    columns=[{"name": i, "id": i} for i in filtered_df.columns],
                    data=[],
                    style_table={'minWidth': '90%', 'overflowY': 'auto', 
                        'border': '1px solid black', 'font-size': '12px'},  # Adjust font size
                    style_cell={'textAlign': 'left', 'padding': '5px', 'border': '1px solid black'},  # Add padding
                    page_action='none'
                ),
            ]),
            width=6
        ),
        dbc.Col(
            html.Div([
                html.H2("Winter 23-24 peaks versus historic levels", style={"textAlign": "center", 'font-size': '14px'}),
                html.Div(id="historic-graph", className="card"),
            ]),
            width=6
        )
    ]),
    html.Hr(),  # line break
    # Top ten table
    dbc.Row([
        dbc.Col(
            html.Div([
                html.H2("These stations experienced exceptional (top-10 all-time) historic levels this winter", style={"textAlign": "center", 'font-size': '14px'}),
                html.Div([
                    dash_table.DataTable(
                        id='top-ten-rank',
                        columns=[{"name": i, "id": i} for i in top_ten_records.columns],
                        data=top_ten_records.to_dict('records'),
                        fixed_columns={'headers': True, 'data': 3},
                        style_table={'minWidth': '100%', 'height': '200px', 'overflowY': 'auto', 
                                     'border': '1px solid black', 'font-size': '12px'},
                        style_cell={'textAlign': 'left', 'padding': '5px', 'border': '1px solid black'},
                        page_action='none',
                        filter_action="native",
                        sort_action="native",
                    ),
                ]),
            ]),
            width=12
        )
    ]),
    html.Hr(),  # line break
    # All peak table
    dbc.Row([
        dbc.Col(
            html.Div([
                html.H2("Here are the peaks for all stations in all events this winter", style={"textAlign": "center", 'font-size': '14px'}),
                dash_table.DataTable(
                    id='peak-table-all',
                    columns=[{"name": i, "id": i} for i in peak_table_all.columns],
                    data=peak_table_all.to_dict('records'),
                    fixed_columns={'headers': True, 'data': 3},
                    style_table={'minWidth': '100%', 'height': '200px', 'overflowY': 'auto', 
                                'border': '1px solid black', 'font-size': '12px'},
                    style_cell={'textAlign': 'left', 'padding': '5px', 'border': '1px solid black'},
                    page_action='none',
                    filter_action="native",
                    sort_action="native"
                ),
            ]),
            width=12
        )
    ]),

    # Download all peak table
    dbc.Row([
        dbc.Col(
            [
                dbc.Button(id='btn',
                    children="Download this data table with peaks for all stations here",
                    color="info",
                    className="mt-1"
                ),
                dcc.Download(id="download-component")
            ],
            width=12
        )
    ]),
    html.Hr(),  # line break
    
    # Storm Parameters used
    html.H4("Data sources used in this app", style={"text-decoration": "underline"}, id="data-info"),
    html.H5("River level data"),
    html.P([" River level data is accessed using the ",
             html.A("DEFRA Hydrology API", href="https://environment.data.gov.uk/hydrology/doc/reference"),
             f". ID codes for a list of selected key gauging stations across the West Midlands was used,\
                  and 15-minute resolution river level data was downloaded for analysis.\
                Data was downloaded for the period {MIN_DATE.strftime('%d %b %Y')} to {MAX_DATE.strftime('%d %b %Y')}"
            ]),
    html.H5("Storm Parameters", style={"textAlign":"left"}),
    html.P("Storm parameters were set to cover time periods consistent with internal EA reporting."),
    html.Div(generate_storm_info()),
    html.H5("Historic data", style={"textAlign":"left"}),
    html.P("Historic data has been digitised from internal EA records (Gaugeboards in SHWG and Thermometer Sheets in SWWM)."),

    html.Hr(),  # line break
    # Final gif
    dbc.Row([
        dbc.Col(
            dbc.CardImg(
                src="https://media.giphy.com/media/v1.Y2lkPTc5MGI3NjExYjRkczhwNHI2emxnMzUzMDBwYWdiN2JpY3k0aGswb3A3cXhiMW1leiZlcD12MV9pbnRlcm5hbF9naWZfYnlfaWQmY3Q9Zw/25Pke1HBWFhjLNrMUe/giphy.gif",
                bottom=True,
                style={"width": "150px", "height": "150px"}
            ),
            width=2
        ),
        dbc.Col(
            html.Div([
                html.P("Yay my code is doing what I want it to."),
                html.P("This makes me a happy Stacy."),
            ]),
            width=10
        )
    ]),
], fluid=True)


### DEFINE CALLBACKS
@app.callback(
    Output("download-component", "data"),
    Input("btn", "n_clicks"),
    prevent_initial_call=True,
)

def func(n_clicks):
    return dcc.send_data_frame(peak_table_all.to_csv, "BEST_DATA_TABLE_EVER.csv")
 
def update_station_options(selected_river):
    if selected_river:
        stations = [key for key, value in data_dict.items() if value['river_name'] == selected_river]
        options = [{'label': station, 'value': station} for station in stations]
        return options
    else:
        return []
    
# Callback associated with the dropdowns
@app.callback(
    [Output('output-graph', 'children'),
     Output('peak-table', 'children'),
     Output('historic-graph', 'children'),
     Output('station-top-ten-table', 'data'),
     Output('station-dropdown', 'options')],
    [Input('river-dropdown', 'value'),
    Input('station-dropdown', 'value')]
)

### UPDATE FUNCTIONS BASED ON THE CALLBACK
def update_graph_peak_table_top_ten(selected_river, selected_station):
    # Update the options of the station dropdown based on the selected river
    station_options = update_station_options(selected_river)

    if selected_station:
        station_data = data_dict[selected_station]
        df = station_data.get('date_values')
        river_name = station_data.get('river_name', 'Unknown River')
        if df is not None:
            figure = {
                'data': [{'x': df['dateTime'], 
                          'y': df['value'], 
                          'type': 'line', 
                          'name': 'River Level'}],
                'layout': {'title': f'River Levels for {selected_station} ({river_name})', 'xaxis': {'title': 'Date Time'},
                           'yaxis': {'title': 'Value'}}
            }
            if selected_station in max_values:
                # Create table rows for peak values
                table_rows = []
                for filter_name, max_value_info in max_values[selected_station].items():
                    max_datetime = max_value_info['dateTime']
                    max_value = max_value_info['value']
                    figure['data'].append({
                        'x': [max_datetime], 
                        'y': [max_value], 
                        'mode': 'markers',
                        'marker': {'color': DATE_FILTERS[filter_name][2], 'size': 10},
                        'name': f'Storm {filter_name} peak'})
                    max_datetime = max_value_info['dateTime'].strftime('%d-%b-%Y %H:%M')  # Format datetime here,after the figure plot,and not at the data entry otherwise peaks don't plot
                    color = DATE_FILTERS[filter_name][2] # get color info

                    # Create a colored marker cell
                    colored_marker_cell = html.Td(style={'background-color': color, 'width': '10px'}, children=[])

                    table_rows.append(html.Tr([
                        colored_marker_cell,  # Colored marker cell
                        html.Td(filter_name, style={'font-size': '16px'}),
                        html.Td(str(max_datetime), style={'font-size': '16px'}),
                        html.Td(str(max_value), style={'font-size': '16px'})
                    ]))

                # Create the table
                peak_table = html.Table([
                    html.Thead(html.Tr([
                        html.Th('', style={'font-size': '16px', 'width': '18px'}),  # Empty header for marker cell
                        html.Th('Storm Name', style={'font-size': '16px'}),
                        html.Th('Date Time', style={'font-size': '16px'}),
                        html.Th('Peak Value', style={'font-size': '16px'})
                    ])),
                    html.Tbody(table_rows)
                ])

                # Get top 10 table data
                top_10_df = station_top_ten(selected_station)

                # Make historic plot
                fig = plot_historic_levels(filtered_df, selected_station, threshold_dict)
                # Include the generated plot in the layout
                historic_graph = dcc.Graph(id='historic-levels-graph', figure=fig)

                # Return graph, peak table, and top 10 table data
                return dcc.Graph(id='river-level-graph', figure=figure), peak_table, historic_graph, top_10_df.to_dict('records'), station_options

    return "No data available for selected station.", "", [], [], station_options

# Callback to update map based on selected station
@app.callback(
    Output('map-container', 'srcDoc'),
    [Input('station-dropdown', 'value')]
)

def update_map(selected_station):
    return create_map(data_dict, selected_station)

# Define callback to toggle the modal
@app.callback(
    dash.dependencies.Output("modal", "is_open"),
    [dash.dependencies.Input("open_modal", "n_clicks"), dash.dependencies.Input("close_modal", "n_clicks")],
    [dash.dependencies.State("modal", "is_open")],
)
def toggle_modal(n1, n2, is_open):
    if n1 or n2:
        return not is_open
    return is_open


# Run the app
if __name__ == '__main__':
    app.run_server(debug=True)